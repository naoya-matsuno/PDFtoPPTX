import fitz
from pptx import Presentation
from pptx.util import Pt
import os
import tkinter
from tkinter import messagebox, filedialog

class Slides():
    __slides: list
    __path: str
    
    def __init__(self, path: str) -> None:
        self.__path, _ = os.path.splitext(path)
        self.__slides = fitz.open(path)

    def convertPPTX(self) -> None:
        dpi = 100
        base = self.__slides[0].get_pixmap(dpi=dpi).width
        if base < self.__slides[0].get_pixmap(dpi=dpi).height:
            base = self.__slides[0].get_pixmap(dpi=dpi).height
        if base < 2000:
            dpi = int(2000 / base * dpi)
        presentation = Presentation()
        presentation.slide_height = Pt(self.__slides[0].get_pixmap(dpi=dpi).height)
        presentation.slide_width = Pt(self.__slides[0].get_pixmap(dpi=dpi).width)
        path = ""
        for slide in self.__slides:
            pptx_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            path = self.__path + "_page.png"
            slide.get_pixmap(dpi=dpi).save(path)
            pptx_slide.shapes.add_picture(path, 0, 0, presentation.slide_width, presentation.slide_height)
        os.remove(path)
        presentation.save(self.__path + "_converted.pptx")

class Window():
    __window: tkinter.Tk
    __title = "PDFtoPPTX"
    __width: int
    __height: int
    __convert_btn: tkinter.Button
    __select_btn: tkinter.Button
    __path = ""
    __path_label: tkinter.Label
    __process_label: tkinter.Label
    
    def __init__(self) -> None:
        self.__window = tkinter.Tk()
        self.__window.title(self.__title)
        self.__width = 600
        self.__height = 200
        window_pos_left = int((self.__window.winfo_screenwidth() - self.__width)/ 2)
        window_pos_top = int((self.__window.winfo_screenheight() - self.__height)/ 2)
        self.__window.geometry(f"{self.__width}x{self.__height}+{window_pos_left}+{window_pos_top}")
        
        self.__path_label = tkinter.Label(self.__window, justify=tkinter.CENTER)
        self.__process_label = tkinter.Label(self.__window, justify=tkinter.CENTER)
        self.__select_btn = tkinter.Button(self.__window, text="Select", command=self.__setPath, justify=tkinter.CENTER, cursor="hand2")
        self.__convert_btn = tkinter.Button(self.__window, text="Convert", command=self.__pushConvertBtn, justify=tkinter.CENTER,  state=tkinter.DISABLED, cursor="arrow")

        self.__path_label.pack(side=tkinter.TOP, pady=10)
        self.__select_btn.pack(side=tkinter.TOP, pady=5)
        self.__process_label.pack(side=tkinter.TOP, pady=5)
        self.__convert_btn.pack(side=tkinter.BOTTOM, pady=10)
        self.__window.mainloop()

    def __pushConvertBtn(self) -> None:
        writable = True
        if os.path.exists(os.path.splitext(self.__path)[0] + "_converted.pptx"):
            writable = messagebox.askyesno(title=self.__title, message="\"" + os.path.splitext(self.__path)[0] + "_converted.pptx\" already exists.\nDo you want to overwrite it?")

        if writable:
            self.__select_btn["state"] = tkinter.DISABLED
            self.__select_btn["cursor"] = "arrow"
            self.__convert_btn["state"] = tkinter.DISABLED
            self.__convert_btn["cursor"] = "arrow"
            self.__process_label["text"] = "Processing..."
            slides = Slides(self.__path)
            slides.convertPPTX()
            self.__process_label["text"] = ""
            messagebox.showinfo(title=self.__title, message="Finished!")
            self.__select_btn["state"] = tkinter.NORMAL
            self.__select_btn["cursor"] = "hand2"
            self.__convert_btn["state"] = tkinter.NORMAL
            self.__convert_btn["cursor"] = "hand2"

    def __setPath(self) -> None:
        path = filedialog.askopenfilename(filetypes=[("PDF File", "*.pdf")])
        self.__path = path
        self.__path_label["text"] = path
        self.__convert_btn["state"] = tkinter.NORMAL
        self.__convert_btn["cursor"] = "hand2"

def main() -> None:
    Window()

if __name__ == "__main__":
    main()