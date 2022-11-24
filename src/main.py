from turtle import screensize, window_height, window_width
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Pt
import os
import tkinter

class Slides:
    __slides: list
    __load = False
    __path: str
    
    def __init__(self, file: str) -> None:
        self.__path, ext = os.path.splitext(file)
        if (ext == ".pdf"):
            self.__slides = convert_from_path(file)
            print(type(self.__slides))
            self.__load = True
        else:
            print("This file is not PDF file.")

    def toPPTX(self) -> None:
        if (self.__load):
            presentation = Presentation()
            presentation.slide_height = Pt(self.__slides[0].height)
            presentation.slide_width = Pt(self.__slides[0].width)
            path = ""
            for slide in self.__slides:
                pptx_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                path = self.__path + "_page.png"
                slide.save(path, "png")
                pptx_slide.shapes.add_picture(path, 0, 0)
            os.remove(path)
            presentation.save(self.__path + ".pptx")

        else:
            print("Could not load File.")

class Window():
    __window: tkinter.Tk
    __width: int
    __height: int
    __convert_btn: tkinter.Button
    __file = ""
    
    def __init__(self) -> None:
        self.__window = tkinter.Tk()
        self.__window.title("PDFtoPPTX")
        self.__width = 760
        self.__height = 500
        window_pos_left = int((self.__window.winfo_screenwidth() - self.__width)/ 2)
        window_pos_top = int((self.__window.winfo_screenheight() - self.__height)/ 2)
        self.__window.geometry(f"{self.__width}x{self.__height}+{window_pos_left}+{window_pos_top}")
        
        self.__convert_btn = tkinter.Button(self.__window, text="convert", command=self.__pushConvert)

        self.__window.mainloop()

    def __pushConvert(self) -> None:
        if len(self.__file) != 0:
        
        else:
            


def main() -> None:
    a = Window()
    #slides = Slides("./files/test.pdf")
    #slides.toPPTX()

if __name__ == "__main__":
    main()